from typing import Optional, Dict, Any
import re
import io
import aiofiles
from loguru import logger

# Optional imports for document processing
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logger.warning("pdfplumber not available. PDF processing will be limited.")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. DOCX processing will be limited.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPTX processing will be limited.")

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger.warning("pytesseract/PIL not available. OCR will be disabled.")


class TextExtractor:
    """Text extraction from various document formats."""
    
    @staticmethod
    async def extract_from_pdf(file_path: str) -> str:
        """Extract text from PDF file."""
        if not PDFPLUMBER_AVAILABLE:
            raise ValueError("PDF processing not available. Install pdfplumber: pip install pdfplumber")
        try:
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}")
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")
    
    @staticmethod
    async def extract_from_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        if not DOCX_AVAILABLE:
            raise ValueError("DOCX processing not available. Install python-docx: pip install python-docx")
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"DOCX extraction error: {str(e)}")
            raise ValueError(f"Failed to extract text from DOCX: {str(e)}")
    
    @staticmethod
    async def extract_from_txt(file_path: str) -> str:
        """Extract text from TXT file."""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                text = await f.read()
            return text.strip()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                async with aiofiles.open(file_path, 'r', encoding='latin-1') as f:
                    text = await f.read()
                return text.strip()
            except Exception as e:
                logger.error(f"TXT extraction error: {str(e)}")
                raise ValueError(f"Failed to extract text from TXT: {str(e)}")
        except Exception as e:
            logger.error(f"TXT extraction error: {str(e)}")
            raise ValueError(f"Failed to extract text from TXT: {str(e)}")
    
    @staticmethod
    async def extract_from_pptx(file_path: str) -> str:
        """Extract text from PPTX file."""
        if not PPTX_AVAILABLE:
            raise ValueError("PPTX processing not available. Install python-pptx: pip install python-pptx")
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PPTX extraction error: {str(e)}")
            raise ValueError(f"Failed to extract text from PPTX: {str(e)}")
    
    @staticmethod
    async def extract_from_image(file_path: str) -> str:
        """Extract text from image using OCR."""
        if not OCR_AVAILABLE:
            raise ValueError("OCR not available. Install pytesseract and Pillow: pip install pytesseract Pillow")
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.error(f"OCR extraction error: {str(e)}")
            raise ValueError(f"Failed to extract text from image: {str(e)}")
    
    @staticmethod
    async def extract_from_file(file_path: str, file_type: str) -> str:
        """Extract text from file based on type."""
        extractors = {
            'pdf': TextExtractor.extract_from_pdf,
            'docx': TextExtractor.extract_from_docx,
            'txt': TextExtractor.extract_from_txt,
            'pptx': TextExtractor.extract_from_pptx,
        }
        
        extractor = extractors.get(file_type.lower())
        if not extractor:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return await extractor(file_path)
    
    @staticmethod
    def extract_from_url(url: str) -> str:
        """Extract text from URL (placeholder - requires web scraping)."""
        # This would require web scraping implementation
        # For now, return placeholder
        logger.warning(f"URL extraction not implemented: {url}")
        return ""
    
    @staticmethod
    def extract_from_youtube(video_url: str) -> str:
        """Extract transcript from YouTube video (placeholder)."""
        # This would require YouTube transcript API integration
        # For now, return placeholder
        logger.warning(f"YouTube transcript extraction not implemented: {video_url}")
        return ""
